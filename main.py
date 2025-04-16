import streamlit as st
import pandas as pd
import plotly.express as px
from fpdf import FPDF
from PIL import Image
import openai
import re
import os
import requests
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Configura tu clave de API de OpenAI
openai.api_key = ""

# Función para leer archivos .txt
def load_txt(file):
    return file.read().decode("utf-8")

# Función para cargar datos
def load_data(file):
    if file.name.endswith('.csv'):
        return pd.read_csv(file)
    elif file.name.endswith('.xlsx'):
        return pd.read_excel(file)
    else:
        st.error("Formato no soportado. Suba un archivo CSV o Excel.")
        return None

# Llamada 1: Análisis exploratorio de datos con ChatGPT
def analyze_data_summary_from_chatgpt(data, idioma):
    prompt = f"""
    Realiza un análisis descriptivo del siguiente dataset:
    - Resume las columnas, indicando sus significados, medidas, rangos, y media (si corresponde).
    - Incluye información relevante para un informe inicial.

    devuelveme tu resultado en el idioma: {idioma}

    Datos (muestra de las primeras filas):
    {data.head().to_string()}
    """
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response["choices"][0]["message"]["content"]

# Llamada 2: Generar gráficas representativas con ChatGPT
def create_representative_graphs(csv_content, nombre_archivo, info):
    prompt = f"""
Tengo el siguiente archivo con este contenido:
{csv_content}

Por favor, genera un código Python que:
1. Lea el archivo, con el nombre: {nombre_archivo}.
2. Construya gráficos representativos para todas, absolutamente todas, las columnas numéricas.
3. Use matplotlib y pandas para los gráficos.

usa temperatura=1

y ten en cuenta si te hace falta esta información = {info}

Devuelve solo el código Python, sin explicaciones ni texto adicional.
"""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "Eres un asistente experto en análisis de datos y visualización."},
            {"role": "user", "content": prompt}
        ]
    )
    response_content = response['choices'][0]['message']['content']
    match = re.search(r"```python(.*?)```", response_content, re.DOTALL)
    if match:
        return match.group(1).strip()
    else:
        raise ValueError("No se encontró código Python en la respuesta de ChatGPT.")

# Función para ejecutar código desde archivo .py y guardar el gráfico
def execute_python_code_and_save_graph(code, filename="output_graph.png"):
    temp_filename = "plot_data.py"
    with open(temp_filename, 'w') as file:
        file.write(code)
    try:
        # Ejecutar el código generado
        exec(code, {"__name__": "__main__"})
        plt.savefig(filename)  # Guardar el gráfico generado por matplotlib
        plt.close()  # Cerrar la figura
    except Exception as e:
        st.error(f"Error al ejecutar el código generado: {e}")
    finally:
        pass
        #os.remove(temp_filename)  # Eliminar archivo temporal

# Llamada 3: Generar una imagen con ChatGPT (utilizando DALL·E)
def generate_image_from_data(data_description):
    truncated_description = data_description[:900]
    prompt = f"""
    Genera una representación visual artística basada en este resumen de datos:
    {truncated_description}
    """
    try:
        response = openai.Image.create(
            prompt=prompt,
            n=1,
            size="1024x1024"
        )
        image_url = response['data'][0]['url']
        
        # Descargar la imagen
        image_response = requests.get(image_url)
        if image_response.status_code == 200:
            with open("imagen_generada.png", "wb") as f:
                f.write(image_response.content)
            #st.image("imagen_generada.png", caption="Imagen Generada Basada en los Datos", use_column_width=True)
        else:
            st.error("No se pudo descargar la imagen.")
    except openai.error.InvalidRequestError as e:
        st.error(f"Error generando la imagen: {e}")
    except Exception as e:
        st.error(f"Error inesperado: {e}")

# Generar gráficas representativas para Streamlit
def create_representative_graphs_streamlit(data):
    # Seleccionar columnas numéricas
    numeric_columns = data.select_dtypes(include=['float64', 'int64']).columns

    # Crear carpeta "plots" si no existe
    output_folder = "plots"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    graphs = []
    num=0
    for column in numeric_columns:
        # Crear histograma
        fig_hist = px.histogram(data, x=column, title=f'Distribución de {column}')
        graphs.append(fig_hist)
        
        # Guardar el histograma
        hist_path = os.path.join(output_folder, f'graph_{str(num)}.png')
        fig_hist.write_image(hist_path)
        num=num+1

        # Crear boxplot
        fig_box = px.box(data, x=column, title=f'Boxplot de Distribución de {column}')
        graphs.append(fig_box)
        
        # Guardar el boxplot
        box_path = os.path.join(output_folder, f'graph_{str(num)}.png')
        fig_box.write_image(box_path)
        num=num+1

    return graphs


def guardar_archivo_en_local(uploaded_file_to_work, nombre_archivo):
    file_name = uploaded_file_to_work.name
    current_dir = os.getcwd()  # Obtiene el directorio actual donde se ejecuta Streamlit
    
    # Construye la ruta completa para guardar el archivo
    file_path = os.path.join(current_dir, nombre_archivo)
    
    # Guarda el archivo en la ubicación actual
    with open(file_path, "wb") as f:
        f.write(uploaded_file_to_work.getbuffer())
    
    st.write(f"El archivo '{file_name}' se guardó en: {file_path} con el nombre {nombre_archivo}")

# Función para exportar informe a PDF
def export_to_pdf(summary, graphs, generated_graph_path, image_path, filename):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", "B", size=12)

    pdf.cell(200, 10, txt="Informe de Datos", ln=True, align="C")
    pdf.ln(10)

    # Agregar resumen
    pdf.multi_cell(0, 10, summary)
    pdf.ln(10)

    # Agregar gráficos
    for i, graph in enumerate(graphs):
        graph.write_image(f"graph_{i}.png")
        pdf.add_page()
        pdf.image(f"graph_{i}.png", x=10, y=50, w=180)
        os.remove(f"graph_{i}.png")  # Eliminar archivo temporal

    # Agregar gráficos generados dinámicamente
    if os.path.exists(generated_graph_path):
        pdf.add_page()
        pdf.cell(200, 10, txt="Gráfico Generado desde Código Python", ln=True, align="C")
        pdf.ln(10)
        pdf.image(generated_graph_path, x=10, y=50, w=180)
        os.remove(generated_graph_path)  # Eliminar archivo temporal

    # Agregar imagen generada
    if image_path and os.path.exists(image_path):
        pdf.add_page()
        pdf.cell(200, 10, txt="Imagen Generada Basada en los Datos", ln=True, align="C")
        pdf.ln(10)
        pdf.image(image_path, x=10, y=50, w=180)

    pdf.output(filename)
    st.success(f"Informe PDF generado con éxito: {filename}")

# Función para centrar las figuras en las diapositivas
def add_centered_picture(slide, image_path, slide_width, slide_height, max_width=6):
    from PIL import Image as PILImage

    # Obtener dimensiones de la imagen
    img = PILImage.open(image_path)
    img_width, img_height = img.size

    # Escalar la imagen si es más ancha que el máximo permitido
    aspect_ratio = img_height / img_width
    if Inches(max_width) < slide_width:
        width = Inches(max_width)
        height = width * aspect_ratio
    else:
        width = slide_width
        height = width * aspect_ratio

    # Calcular la posición centrada
    left = (slide_width - width) / 2
    top = (slide_height - height) / 2

    # Agregar la imagen centrada
    slide.shapes.add_picture(image_path, left, top, width=width)

# Función para exportar informe a PowerPoint
def export_to_pptx(summary, graphs, generated_graph_path, image_path, filename):
    ppt = Presentation()
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height
    slide_layout = ppt.slide_layouts[5]  # Diseño de diapositiva en blanco

    # Diapositiva inicial con título
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Informe de Datos"

    # Diapositiva para resumen
    slide = ppt.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    textbox.text = "Resumen de Datos:\n" + summary

    # Diapositivas para gráficos
    for i, graph in enumerate(graphs):
        graph_path = f"graph_{i}.png"
        graph.write_image(graph_path)
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Gráfico {i + 1}"
        add_centered_picture(slide, graph_path, slide_width, slide_height)
        os.remove(graph_path)  # Eliminar archivo temporal

    # Diapositiva para gráfico generado dinámicamente
    if os.path.exists(generated_graph_path):
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Gráfico Generado desde Código Python"
        add_centered_picture(slide, generated_graph_path, slide_width, slide_height)

    # Diapositiva para imagen generada
    if image_path and os.path.exists(image_path):
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Imagen Generada Basada en los Datos"
        add_centered_picture(slide, image_path, slide_width, slide_height)

    ppt.save(filename)
    st.success(f"Informe PowerPoint generado con éxito: {filename}")

# Función para exportar informe a HTML
def export_to_html(summary, graphs, generated_graph_path, image_path, filename):
    # Usar un f-string en lugar de .format()
    html_content = f"""
    <!DOCTYPE html>
    <html lang="es">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Informe de Datos</title>
        <style>
            body {{
                font-family: Arial, sans-serif;
                margin: 0;
                padding: 20px;
            }}
            h1, h2 {{
                text-align: center;
            }}
            .section {{
                margin-bottom: 40px;
            }}
            img {{
                max-width: 100%;
                height: auto;
                display: block;
                margin: 0 auto;
            }}
        </style>
    </head>
    <body>
        <h1>Informe de Datos</h1>
        <div class="section">
            <h2>Resumen de Datos</h2>
            <p>{summary.replace('\n', '<br>')}</p> <!-- Reemplazar saltos de línea por <br> -->
        </div>
    """

    # Agregar gráficos
    for i, graph in enumerate(graphs):
        graph_path = f"graph_{i}.png"
        graph.write_image(graph_path)
        html_content += f"""
        <div class="section">
            <h2>Gráfico {i + 1}</h2>
            <img src="plots/{graph_path}" alt="Gráfico {i + 1}">
        </div>
        """
        os.remove(graph_path)  # Eliminar archivo temporal

    # Agregar gráfico generado dinámicamente
    if os.path.exists(generated_graph_path):
        html_content += f"""
        <div class="section">
            <h2>Gráfico Generado desde Código Python</h2>
            <img src="{generated_graph_path}" alt="Gráfico Generado">
        </div>
        """

    # Agregar imagen generada
    if image_path and os.path.exists(image_path):
        html_content += f"""
        <div class="section">
            <h2>Imagen Generada Basada en los Datos</h2>
            <img src="{image_path}" alt="Imagen Generada">
        </div>
        """

    # Cerrar el HTML
    html_content += """
    </body>
    </html>
    """

    # Guardar el archivo HTML
    with open(filename, "w", encoding="utf-8") as file:
        file.write(html_content)
    
    st.success(f"Informe HTML generado con éxito: {filename}")
    st.markdown(f"[Descargar el Informe HTML](./{filename})", unsafe_allow_html=True)



### ----- Main ---------------------------

# Interfaz de Streamlit
st.title("Generador Automático de Informes con IA")

# inicialización de variables:
description_file_presente = False
resumen_file_presente = False

language = st.selectbox("Selecciona el idioma", options=["Español", "English"])


# Pestaña para cargar archivo description
description_file = st.file_uploader("Sube tu archivo description (.txt)", type=["txt"])
if description_file:
    description = load_txt(description_file)
    st.subheader("Contenido de Description")
    st.text_area("Contenido cargado:", description, height=300)
    guardar_archivo_en_local(description_file, "description_file.txt")
    description_file_presente = True

# Pestaña para cargar archivo resumen
resumen_file = st.file_uploader("Sube tu archivo resumen (.txt)", type=["txt"])
if resumen_file:
    resumen = load_txt(resumen_file)
    st.subheader("Contenido de Resumen")
    st.text_area("Contenido cargado:", resumen, height=300)
    guardar_archivo_en_local(resumen_file, "resumen_file.txt")
    resumen_file_presente = True

uploaded_file = st.file_uploader("Sube tu archivo de datos (CSV o Excel)", type=["csv", "xlsx"])

if uploaded_file:
    data = load_data(uploaded_file)

    guardar_archivo_en_local(uploaded_file, uploaded_file.name)

    if data is not None:
        st.subheader("Vista Previa de los Datos")
        st.dataframe(data)

        if resumen_file_presente:
            data_summary = resumen
        else:
            st.subheader("Generando Resumen de Datos...")
            data_summary = analyze_data_summary_from_chatgpt(data, language)
            st.text_area("Resumen del Dataset", data_summary, height=300)

        st.subheader("Generando Gráficas de cada variable en...")
        graphs_streamlit = create_representative_graphs_streamlit(data)
        for graph in graphs_streamlit:
            st.plotly_chart(graph)

        st.subheader("Generando Código para Gráficas...")
        csv_content = data.head().to_csv(index=False)
        graphs_code = create_representative_graphs(csv_content, uploaded_file.name, info=data_summary)
        st.code(graphs_code, language="python")

        st.subheader("Ejecutando Código y Generando Gráfico...")
        generated_graph_path = "output_graph.png"
        execute_python_code_and_save_graph(graphs_code, generated_graph_path)
        st.image(generated_graph_path, caption="Gráfico Generado por Código", use_column_width=True)

        st.subheader("Generando Imagen...")
        generate_image_from_data(data_summary)
        #st.image("imagen_generada.png", caption="Imagen Generada Basada en los Datos", use_column_width=True)
        #image_path="perrete.jpg"
        image_path="imagen_generada.png"
        st.image(image_path, caption="Imagen Generada Basada en los Datos", use_column_width=True)

        if st.button("Generar Informe en PDF"):
            export_to_pdf(data_summary, graphs_streamlit, generated_graph_path, image_path, "informe.pdf")

        if st.button("Generar Informe en PowerPoint"):
            export_to_pptx(data_summary, graphs_streamlit, generated_graph_path, image_path, "informe.pptx")

        # Agregar botón para generar informe en HTML
        if st.button("Generar Informe en HTML"):
            export_to_html(data_summary, graphs_streamlit, generated_graph_path, image_path, "informe.html")
    
